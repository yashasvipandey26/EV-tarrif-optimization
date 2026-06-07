from __future__ import annotations

from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ASSET_DIR = Path("ppt_assets")
ASSET_DIR.mkdir(exist_ok=True)
DECK_PATH = Path("EV_Dynamic_Tariff_Optimization_Deck.pptx")
OUTLINE_PATH = Path("EV_Dynamic_Tariff_Optimization_Deck_Outline.md")

COLORS = {
    "navy": RGBColor(17, 34, 64),
    "blue": RGBColor(39, 116, 174),
    "cyan": RGBColor(74, 193, 210),
    "green": RGBColor(42, 157, 143),
    "orange": RGBColor(244, 162, 97),
    "red": RGBColor(231, 111, 81),
    "gray": RGBColor(92, 99, 112),
    "light": RGBColor(245, 247, 250),
    "white": RGBColor(255, 255, 255),
}
MPL = {
    "navy": "#112240",
    "blue": "#2774AE",
    "cyan": "#4AC1D2",
    "green": "#2A9D8F",
    "orange": "#F4A261",
    "red": "#E76F51",
    "gray": "#5C6370",
    "light": "#F5F7FA",
}

plt.rcParams.update({
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "axes.edgecolor": "#CBD5E1",
    "axes.labelcolor": MPL["gray"],
    "xtick.color": MPL["gray"],
    "ytick.color": MPL["gray"],
    "font.size": 10,
})


def load_data():
    metrics = pd.read_csv("model_evaluation_summary.csv")
    tariff = pd.read_csv("dynamic_tariff_table.csv", parse_dates=["hour_start"])
    features = pd.read_csv("engineered_station_hourly_features.csv", parse_dates=["hour_start"])
    stations = pd.read_csv("station_utilization_indices.csv")
    sim = pd.read_csv("monitoring_learning_simulation.csv", parse_dates=["hour_start"])
    return metrics, tariff, features, stations, sim


def metric(metrics: pd.DataFrame, name: str) -> float:
    return float(metrics.loc[metrics["metric"].eq(name), "value"].iloc[0])


def savefig(path: Path):
    plt.tight_layout()
    plt.savefig(path, dpi=220, bbox_inches="tight")
    plt.close()


def make_charts(metrics, tariff, features, stations, sim):
    paths = {}

    demand = features.copy()
    demand["hour"] = demand["hour_start"].dt.hour
    demand["day"] = demand["hour_start"].dt.day_name().str[:3]
    order = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    heat = demand.pivot_table(index="day", columns="hour", values="volume_kwh", aggfunc="mean").reindex(order)
    fig, ax = plt.subplots(figsize=(8.2, 3.9))
    im = ax.imshow(heat.fillna(0), aspect="auto", cmap="YlGnBu")
    ax.set_title("Average Charging Load by Day and Hour", loc="left", fontweight="bold", color=MPL["navy"])
    ax.set_xlabel("Hour of day")
    ax.set_ylabel("Day")
    ax.set_xticks(range(0, 24, 2))
    ax.set_xticklabels(range(0, 24, 2))
    ax.set_yticks(range(len(order)))
    ax.set_yticklabels(order)
    cbar = fig.colorbar(im, ax=ax, fraction=0.03, pad=0.02)
    cbar.set_label("kWh")
    paths["heatmap"] = ASSET_DIR / "demand_heatmap.png"
    savefig(paths["heatmap"])

    hourly = demand.groupby("hour", observed=True).agg(
        avg_volume=("volume_kwh", "mean"), avg_util=("utilization_rate", "mean"), avg_queue=("queue_length_proxy", "mean")
    ).reset_index()
    fig, ax1 = plt.subplots(figsize=(8.2, 3.9))
    ax1.plot(hourly["hour"], hourly["avg_volume"], color=MPL["blue"], linewidth=2.5, label="Avg kWh")
    ax1.set_ylabel("Average kWh", color=MPL["blue"])
    ax2 = ax1.twinx()
    ax2.plot(hourly["hour"], hourly["avg_util"], color=MPL["orange"], linewidth=2.5, label="Utilization")
    ax2.axhline(0.8, color=MPL["red"], linestyle="--", linewidth=1.5, label="Surge threshold")
    ax2.axhline(0.3, color=MPL["green"], linestyle="--", linewidth=1.5, label="Discount threshold")
    ax2.set_ylabel("Utilization rate", color=MPL["orange"])
    ax1.set_title("Hourly Demand and Tariff Trigger Context", loc="left", fontweight="bold", color=MPL["navy"])
    ax1.set_xlabel("Hour of day")
    ax1.set_xticks(range(0, 24, 2))
    lines = ax1.get_lines() + ax2.get_lines()
    ax1.legend(lines, [l.get_label() for l in lines], loc="upper left", frameon=False, ncol=2)
    paths["hourly"] = ASSET_DIR / "hourly_demand_utilization.png"
    savefig(paths["hourly"])

    sample = tariff.sample(min(5000, len(tariff)), random_state=42).sort_values("target_volume_kwh")
    fig, ax = plt.subplots(figsize=(7.4, 4.0))
    ax.scatter(sample["target_volume_kwh"], sample["predicted_volume_kwh"], s=7, alpha=0.25, color=MPL["blue"])
    low = min(sample["target_volume_kwh"].min(), sample["predicted_volume_kwh"].min())
    high = max(sample["target_volume_kwh"].max(), sample["predicted_volume_kwh"].max())
    ax.plot([low, high], [low, high], color=MPL["red"], linestyle="--", linewidth=1.5)
    ax.set_title("Demand Model: Actual vs Predicted Hour-Ahead Load", loc="left", fontweight="bold", color=MPL["navy"])
    ax.set_xlabel("Actual next-hour kWh")
    ax.set_ylabel("Predicted next-hour kWh")
    ax.text(0.03, 0.94, f"R²={metric(metrics, 'R2 Score'):.3f}\nMAE={metric(metrics, 'MAE'):.1f} kWh", transform=ax.transAxes,
            bbox=dict(boxstyle="round,pad=0.35", facecolor="white", edgecolor="#CBD5E1"), va="top")
    paths["actual_pred"] = ASSET_DIR / "actual_vs_predicted.png"
    savefig(paths["actual_pred"])

    signal_counts = tariff["pricing_signal"].value_counts().reindex(["DISCOUNT", "BASELINE", "SURGE"]).fillna(0)
    fig, axes = plt.subplots(1, 2, figsize=(8.6, 3.9), gridspec_kw={"width_ratios": [1, 1.1]})
    colors = [MPL["green"], MPL["gray"], MPL["red"]]
    axes[0].bar(signal_counts.index, signal_counts.values / signal_counts.sum() * 100, color=colors)
    axes[0].set_title("Pricing Signal Mix", loc="left", fontweight="bold", color=MPL["navy"])
    axes[0].set_ylabel("% station-hours")
    axes[0].tick_params(axis="x", rotation=20)
    axes[1].hist(tariff["dynamic_tariff_rupee_per_kwh"], bins=32, color=MPL["blue"], alpha=0.85)
    axes[1].axvline(15, color=MPL["orange"], linestyle="--", linewidth=2, label="₹15 baseline")
    axes[1].set_title("Tariff Distribution", loc="left", fontweight="bold", color=MPL["navy"])
    axes[1].set_xlabel("₹ / kWh")
    axes[1].legend(frameon=False)
    paths["pricing"] = ASSET_DIR / "pricing_outcomes.png"
    savefig(paths["pricing"])

    fig, axes = plt.subplots(1, 3, figsize=(9.2, 3.6))
    rev_gain = metric(metrics, "Revenue Gain %")
    wait_red = metric(metrics, "Average Waiting Time Reduction %")
    off_peak = metric(metrics, "Off-Peak Uplift %")
    for ax, title, value, color in zip(
        axes,
        ["Revenue Gain", "Waiting Proxy Reduction", "Off-Peak Uplift"],
        [rev_gain, wait_red, off_peak],
        [MPL["green"], MPL["blue"], MPL["orange"]],
    ):
        ax.bar([title], [value], color=color, width=0.55)
        ax.set_ylim(0, max(45, value * 1.25))
        ax.set_ylabel("%")
        ax.text(0, value + 1, f"{value:.1f}%", ha="center", fontweight="bold", color=MPL["navy"])
        ax.set_xticklabels([title], rotation=18, ha="right")
    fig.suptitle("Monitoring Agent Simulation: Directional Operational Outcomes", fontweight="bold", color=MPL["navy"], x=0.02, ha="left")
    paths["monitoring"] = ASSET_DIR / "monitoring_outcomes.png"
    savefig(paths["monitoring"])

    top = stations.sort_values("total_volume_kwh", ascending=False).head(12).sort_values("total_volume_kwh")
    fig, ax = plt.subplots(figsize=(8.5, 4.0))
    ax.barh(top["grid"].astype(str), top["total_volume_kwh"] / 1000, color=MPL["cyan"])
    ax.set_title("Top High-Load District Grids", loc="left", fontweight="bold", color=MPL["navy"])
    ax.set_xlabel("Total observed MWh")
    ax.set_ylabel("Grid")
    paths["top_grids"] = ASSET_DIR / "top_grid_loads.png"
    savefig(paths["top_grids"])

    fig, ax = plt.subplots(figsize=(7.4, 3.9))
    ax.scatter(stations["avg_utilization_rate"], stations["total_volume_kwh"] / 1000,
               s=np.clip(stations["charger_count"], 10, 250), alpha=0.55, color=MPL["blue"], edgecolor="white")
    ax.axvline(0.8, color=MPL["red"], linestyle="--", linewidth=1.5)
    ax.axvline(0.3, color=MPL["green"], linestyle="--", linewidth=1.5)
    ax.set_title("Station Portfolio: Utilization vs Energy Throughput", loc="left", fontweight="bold", color=MPL["navy"])
    ax.set_xlabel("Average utilization rate")
    ax.set_ylabel("Total observed MWh")
    paths["portfolio"] = ASSET_DIR / "portfolio_utilization_throughput.png"
    savefig(paths["portfolio"])

    checks = pd.DataFrame({
        "Check": ["Chronological split", "Tariff floor/ceiling", "Spatial imputation", "Elasticity simulation", "Causal caution"],
        "Status": ["Used", "Enforced", "Forward-fill + neighbor mean", "Directional", "No causal claim"],
    })
    fig, ax = plt.subplots(figsize=(8.2, 3.4))
    ax.axis("off")
    table = ax.table(cellText=checks.values, colLabels=checks.columns, cellLoc="left", colLoc="left", loc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 1.55)
    for (row, col), cell in table.get_celld().items():
        cell.set_edgecolor("#CBD5E1")
        if row == 0:
            cell.set_facecolor(MPL["navy"])
            cell.get_text().set_color("white")
            cell.get_text().set_fontweight("bold")
    ax.set_title("Robustness and Transparency Checklist", loc="left", fontweight="bold", color=MPL["navy"])
    paths["checks"] = ASSET_DIR / "robustness_checks.png"
    savefig(paths["checks"])

    return paths


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["navy"]
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.82), Inches(11.8), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(10.5)
        sp.font.color.rgb = COLORS["gray"]


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"EV Dynamic Tariff Optimization | Internal analysis | {page}"
    p.font.size = Pt(8.5)
    p.font.color.rgb = COLORS["gray"]
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, bullets, x, y, w, h, font_size=14, color="navy"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color]
        p.space_after = Pt(8)
    return box


def add_kpi_card(slide, x, y, w, h, label, value, color="blue"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["light"]
    shape.line.color.rgb = COLORS[color]
    shape.line.width = Pt(1.5)
    t = shape.text_frame
    t.clear()
    p = t.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS[color]
    p.alignment = PP_ALIGN.CENTER
    p2 = t.add_paragraph()
    p2.text = label
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLORS["gray"]
    p2.alignment = PP_ALIGN.CENTER


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck(metrics, tariff, features, stations, sim, paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    page = 0
    slide = blank_slide(prs)
    page += 1
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["navy"]
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(8.2), Inches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "Agentic AI-Based Dynamic Tariff Optimization"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(2.72), Inches(8.8), Inches(0.7))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "ACN session logs + UrbanEV Shenzhen district panel | ₹15/kWh baseline comparison"
    sp.font.size = Pt(16)
    sp.font.color.rgb = COLORS["cyan"]
    add_kpi_card(slide, 0.85, 4.25, 2.2, 0.9, "Revenue gain", f"{metric(metrics, 'Revenue Gain %'):.1f}%", "green")
    add_kpi_card(slide, 3.25, 4.25, 2.2, 0.9, "Model R²", f"{metric(metrics, 'R2 Score'):.3f}", "cyan")
    add_kpi_card(slide, 5.65, 4.25, 2.4, 0.9, "Waiting proxy ↓", f"{metric(metrics, 'Average Waiting Time Reduction %'):.1f}%", "orange")
    note = slide.shapes.add_textbox(Inches(0.82), Inches(6.55), Inches(9.5), Inches(0.25))
    npg = note.text_frame.paragraphs[0]
    npg.text = "Prepared from generated pipeline outputs; operational results are simulation-based and should be validated before deployment."
    npg.font.size = Pt(9.5)
    npg.font.color.rgb = COLORS["white"]

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "Executive Summary", "What the analysis supports — and what remains assumption-bound")
    add_kpi_card(slide, 0.55, 1.35, 2.35, 1.0, "Predicted revenue lift", f"{metric(metrics, 'Revenue Gain %'):.1f}%", "green")
    add_kpi_card(slide, 3.1, 1.35, 2.35, 1.0, "Off-peak uplift", f"{metric(metrics, 'Off-Peak Uplift %'):.1f}%", "orange")
    add_kpi_card(slide, 5.65, 1.35, 2.35, 1.0, "Utilization stabilization", f"{metric(metrics, 'Net Utilization Stabilization %'):.1f}%", "blue")
    add_kpi_card(slide, 8.2, 1.35, 2.35, 1.0, "Demand model R²", f"{metric(metrics, 'R2 Score'):.3f}", "cyan")
    add_bullets(slide, [
        "The system combines demand forecasting, tariff decision rules, and simulated feedback to compare dynamic pricing against a ₹15/kWh flat baseline.",
        "Findings suggest targeted discounts can stimulate low-utilization periods while capped surge pricing moderates congestion risk.",
        "Reported gains are directional simulation outcomes, not causal estimates of driver behavior.",
        "Next decision point: pilot limited dynamic tariffs with randomized or phased rollout measurement before policy-wide adoption.",
    ], 0.8, 3.0, 6.3, 2.6, 15)
    slide.shapes.add_picture(str(paths["monitoring"]), Inches(7.35), Inches(3.0), width=Inches(5.3))
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "1. Data Landscape and Preprocessing Decisions", "Granular ACN sessions were aligned to the UrbanEV station-hour panel")
    add_bullets(slide, [
        "UrbanEV: 247 district grids across 720 hourly windows after aggregating 5-minute records from volume, duration, occupancy, and price files.",
        "ACN: vehicle-session logs converted into day/hour profile features for session count, energy, duration, and station activity.",
        "Feature layer includes utilization, revenue, queue proxy, occupancy density, spatial clusters, and neighborhood averages from adjacency/distance matrices.",
        "Missing intervals are handled using forward/back-fill first, then spatial-neighbor average as a transparent operational proxy.",
    ], 0.6, 1.3, 5.7, 4.6, 13.5)
    slide.shapes.add_picture(str(paths["portfolio"]), Inches(6.55), Inches(1.32), width=Inches(6.15))
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "2. EDA Findings and Demand Behavior Insights", "Demand is uneven across time and station portfolio; pricing should avoid one-size-fits-all rules")
    slide.shapes.add_picture(str(paths["heatmap"]), Inches(0.55), Inches(1.2), width=Inches(6.25))
    slide.shapes.add_picture(str(paths["top_grids"]), Inches(6.95), Inches(1.2), width=Inches(5.95))
    add_bullets(slide, [
        "Observed behavior varies by hour/day and by district grid throughput.",
        "Portfolio-level imbalance supports localized tariff decisions rather than a universal citywide multiplier.",
        "Interpretation is associational: patterns describe observed load, not proven causes of demand shifts.",
    ], 0.8, 5.75, 11.7, 0.9, 11.5, "gray")
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "3. Demand Prediction Modeling and Results", "LightGBM forecast produces hour-ahead load, utilization, and congestion probability")
    slide.shapes.add_picture(str(paths["actual_pred"]), Inches(0.6), Inches(1.25), width=Inches(6.1))
    add_kpi_card(slide, 7.05, 1.45, 2.25, 0.85, "RMSE", f"{metric(metrics, 'RMSE'):.1f}", "blue")
    add_kpi_card(slide, 9.55, 1.45, 2.25, 0.85, "MAE", f"{metric(metrics, 'MAE'):.1f}", "orange")
    add_kpi_card(slide, 7.05, 2.55, 2.25, 0.85, "R² score", f"{metric(metrics, 'R2 Score'):.3f}", "green")
    add_kpi_card(slide, 9.55, 2.55, 2.25, 0.85, "Util. RMSE", f"{metric(metrics, 'Utilization RMSE'):.3f}", "cyan")
    add_bullets(slide, [
        "Features: temporal cycles, charger capacity, occupancy density, queue proxy, spatial clusters, and neighboring-grid context.",
        "Chronological split is used to mimic forward-looking deployment more closely than random station-hour mixing.",
        "Congestion probability is calibrated through a logistic proxy around the 80% utilization threshold.",
    ], 7.05, 4.0, 5.45, 1.65, 12.2)
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "4. Dynamic Tariff Optimization Logic and Pricing Outcomes", "Rule-constrained pricing maps forecasted utilization to bounded tariff actions")
    slide.shapes.add_picture(str(paths["hourly"]), Inches(0.55), Inches(1.2), width=Inches(6.15))
    slide.shapes.add_picture(str(paths["pricing"]), Inches(6.85), Inches(1.2), width=Inches(6.0))
    add_bullets(slide, [
        "Surge trigger: predicted utilization >80%; discount trigger: predicted utilization <30%.",
        "Tariffs are bounded between 0.85x and 1.45x of ₹15/kWh to prevent irrational compounding.",
        "Pricing outputs are recommendations for operational testing, not guarantees of realized driver response.",
    ], 0.85, 5.82, 11.7, 0.75, 11.5, "gray")
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "5. Monitoring Agent Evaluation and Feedback Loop", "Simulated elasticity links prices to load, congestion, revenue, and waiting proxy")
    slide.shapes.add_picture(str(paths["monitoring"]), Inches(0.65), Inches(1.25), width=Inches(7.2))
    add_bullets(slide, [
        "Feedback loop applies a negative elasticity assumption: higher prices reduce modeled demand; discounts increase off-peak uptake.",
        "Simulation indicates +5.3% revenue over flat pricing, +18.9% off-peak uplift, and 38.6% lower waiting-time proxy.",
        "These are scenario estimates; pilot telemetry should recalibrate elasticity and customer acceptance parameters.",
    ], 8.15, 1.45, 4.5, 3.6, 14)
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "6. Business, Operational, and Policy Implications", "A controlled pilot can test whether simulated gains translate into field outcomes")
    add_bullets(slide, [
        "Business: dynamic tariffs may improve revenue efficiency while preserving bounded customer-facing prices.",
        "Operations: station-level signals can smooth overload risk and redirect demand to under-used windows.",
        "Policy: transparent thresholds, tariff caps, and published discount windows can reduce fairness concerns.",
        "Governance: monitor adverse impacts by geography, station class, and time-of-day before scaling.",
        "Recommended next step: A/B or phased rollout with opt-out safeguards and weekly model recalibration.",
    ], 0.75, 1.35, 6.1, 4.6, 15)
    add_kpi_card(slide, 7.45, 1.45, 2.35, 1.0, "Efficiency score", f"₹{metric(metrics, 'Pricing Efficiency Score'):.2f}/kWh", "green")
    add_kpi_card(slide, 10.15, 1.45, 2.35, 1.0, "Baseline", "₹15/kWh", "orange")
    add_bullets(slide, [
        "Avoid causal language until randomized or phased pilots measure observed demand response.",
        "Treat elasticity, queue proxy, and imputation as explicit model assumptions.",
        "Use guardrails: cap surge, preserve discounts, audit service equity.",
    ], 7.45, 3.05, 4.9, 2.25, 13.5, "gray")
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "Appendix A. Additional Analysis and Robustness Checks", "Transparency checklist for the current analytical build")
    slide.shapes.add_picture(str(paths["checks"]), Inches(0.75), Inches(1.25), width=Inches(6.2))
    add_bullets(slide, [
        "Model family: LightGBM if available; sklearn gradient boosting fallback is coded for portability.",
        "Core target: next-hour station-grid kWh; secondary target: next-hour utilization rate.",
        "Robustness extension: test tariff results under low/medium/high elasticity scenarios before launch.",
        "Validation extension: compare against naive seasonal baseline and rolling-origin backtests.",
    ], 7.25, 1.4, 5.2, 3.4, 13.5)
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "Appendix B. Assumptions and Limitations", "Current findings support decisioning hypotheses, not definitive causal estimates")
    add_bullets(slide, [
        "ACN and UrbanEV represent different geographies and infrastructures; ACN is used as cross-dataset behavioral signal, not direct Shenzhen ground truth.",
        "Queue length and waiting time are proxies derived from occupancy/utilization rather than observed physical queues.",
        "Price elasticity is simulated from an assumed curve; real driver response must be measured in production pilots.",
        "Spatial neighbor imputation assumes adjacent grids are operationally similar enough to proxy missing station-hour values.",
        "Revenue lift depends on tariff acceptance, regulatory constraints, app communication, and competitive alternatives.",
    ], 0.9, 1.35, 11.2, 4.8, 15)
    add_footer(slide, page)

    slide = blank_slide(prs)
    page += 1
    add_title(slide, "Appendix C. Output Artifacts", "Files generated in the workspace for review and reuse")
    add_bullets(slide, [
        "ev_tariff_optimization_pipeline.py — reproducible data, model, pricing, and monitoring pipeline.",
        "model_evaluation_summary.csv — final benchmark metrics for prediction, pricing, and system quality.",
        "dynamic_tariff_table.csv — station-hour tariff recommendations and expected revenue fields.",
        "station_utilization_indices.csv — grid-level utilization, queue, throughput, and cluster summaries.",
        "monitoring_learning_simulation.csv — simulated feedback-cycle outcomes on the test horizon.",
        "engineered_station_hourly_features.csv — compact station-hour feature table for EDA and audit.",
    ], 0.9, 1.35, 11.2, 4.8, 14.5)
    add_footer(slide, page)

    prs.save(DECK_PATH)


def write_outline(metrics):
    outline = f"""# EV Dynamic Tariff Optimization Deck Outline

## Cover
- Agentic AI-Based Dynamic Tariff Optimization
- Datasets: ACN session logs and UrbanEV Shenzhen district panel
- Baseline: flat ₹15/kWh pricing

## Executive Summary
- Forecasting + pricing + monitoring loop suggests directional operational upside.
- Key metrics: Revenue gain {metric(metrics, 'Revenue Gain %'):.1f}%, R2 {metric(metrics, 'R2 Score'):.3f}, off-peak uplift {metric(metrics, 'Off-Peak Uplift %'):.1f}%, waiting proxy reduction {metric(metrics, 'Average Waiting Time Reduction %'):.1f}%.
- Caveat: simulation results are not causal claims.

## Main Slides
1. Data landscape and preprocessing decisions
2. EDA findings and demand behavior insights
3. Demand prediction modeling and results
4. Dynamic tariff optimization logic and pricing outcomes
5. Monitoring agent evaluation and feedback loop performance
6. Business, operational, and policy implications

## Appendix
- Robustness and transparency checks
- Assumptions and limitations
- Output artifacts generated in workspace
"""
    OUTLINE_PATH.write_text(outline)


def main():
    metrics, tariff, features, stations, sim = load_data()
    paths = make_charts(metrics, tariff, features, stations, sim)
    build_deck(metrics, tariff, features, stations, sim, paths)
    write_outline(metrics)
    print(f"Saved {DECK_PATH}")
    print(f"Saved {OUTLINE_PATH}")
    print(f"Saved charts in {ASSET_DIR}/")


if __name__ == "__main__":
    main()
